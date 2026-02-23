"""
Cortex AI File Parser
Handles extraction of text content from any file type.
Supports PDF, DOCX, XLSX, CSV, JSON, Markdown, HTML, code files, and more.
Falls back to raw text extraction for unknown formats.
Created by Geo Cherian Mathew.
"""
import os
import csv
import io
import json
import re
from pathlib import Path
from typing import Optional
from dataclasses import dataclass, field
from datetime import datetime

import pypdf
import docx
import pandas as pd


@dataclass
class ParsedDocument:
    """Represents a parsed document with its content and metadata."""
    filename: str
    file_type: str
    content: str
    page_count: int = 1
    word_count: int = 0
    char_count: int = 0
    metadata: dict = field(default_factory=dict)
    parsed_at: str = field(default_factory=lambda: datetime.utcnow().isoformat())

    def __post_init__(self):
        self.word_count = len(self.content.split())
        self.char_count = len(self.content)


# ═══════════════════════════════════════════════════════════════════════
#  SPECIALIZED PARSERS
# ═══════════════════════════════════════════════════════════════════════

def parse_pdf(file_path: str) -> ParsedDocument:
    """Extract text from a PDF file."""
    reader = pypdf.PdfReader(file_path)
    pages = []
    for i, page in enumerate(reader.pages):
        text = page.extract_text() or ""
        if text.strip():
            pages.append(f"[Page {i + 1}]\n{text}")

    content = "\n\n".join(pages)
    return ParsedDocument(
        filename=Path(file_path).name,
        file_type="pdf",
        content=content,
        page_count=len(reader.pages),
        metadata={
            "pdf_info": {k: str(v) for k, v in (reader.metadata or {}).items()},
            "total_pages": len(reader.pages),
        },
    )


def parse_docx(file_path: str) -> ParsedDocument:
    """Extract text from a DOCX file."""
    doc = docx.Document(file_path)
    paragraphs = []
    for para in doc.paragraphs:
        if para.text.strip():
            paragraphs.append(para.text)

    # Also extract tables
    tables_text = []
    for i, table in enumerate(doc.tables):
        table_rows = []
        for row in table.rows:
            cells = [cell.text.strip() for cell in row.cells]
            table_rows.append(" | ".join(cells))
        tables_text.append(f"\n[Table {i + 1}]\n" + "\n".join(table_rows))

    content = "\n".join(paragraphs)
    if tables_text:
        content += "\n\n--- Tables ---\n" + "\n".join(tables_text)

    return ParsedDocument(
        filename=Path(file_path).name,
        file_type="docx",
        content=content,
        metadata={"paragraph_count": len(paragraphs), "table_count": len(doc.tables)},
    )


def parse_csv(file_path: str) -> ParsedDocument:
    """Extract text from a CSV file. Converts to readable structured format."""
    df = pd.read_csv(file_path)
    
    # Build a readable representation
    lines = []
    lines.append(f"CSV File with {len(df)} rows and {len(df.columns)} columns.")
    lines.append(f"Columns: {', '.join(df.columns.tolist())}")
    lines.append("")

    # Summary statistics for numeric columns
    numeric_cols = df.select_dtypes(include=["number"]).columns.tolist()
    if numeric_cols:
        lines.append("--- Numeric Summary ---")
        for col in numeric_cols:
            lines.append(
                f"  {col}: min={df[col].min()}, max={df[col].max()}, "
                f"mean={df[col].mean():.2f}, median={df[col].median():.2f}"
            )
        lines.append("")

    # Full data as markdown-style table (limit to 200 rows for very large files)
    display_df = df.head(200) if len(df) > 200 else df
    lines.append("--- Data ---")
    lines.append(display_df.to_string(index=False))

    if len(df) > 200:
        lines.append(f"\n... and {len(df) - 200} more rows (truncated)")

    content = "\n".join(lines)
    return ParsedDocument(
        filename=Path(file_path).name,
        file_type="csv",
        content=content,
        metadata={
            "row_count": len(df),
            "column_count": len(df.columns),
            "columns": df.columns.tolist(),
        },
    )


def parse_excel(file_path: str) -> ParsedDocument:
    """Extract text from Excel files (.xlsx, .xls)."""
    try:
        xls = pd.ExcelFile(file_path)
        all_content = []
        total_rows = 0
        
        for sheet_name in xls.sheet_names:
            df = pd.read_excel(xls, sheet_name=sheet_name)
            total_rows += len(df)
            
            lines = []
            lines.append(f"\n=== Sheet: {sheet_name} ({len(df)} rows × {len(df.columns)} columns) ===")
            lines.append(f"Columns: {', '.join(df.columns.astype(str).tolist())}")
            
            # Numeric summary
            numeric_cols = df.select_dtypes(include=["number"]).columns.tolist()
            if numeric_cols:
                lines.append("\n--- Numeric Summary ---")
                for col in numeric_cols:
                    lines.append(
                        f"  {col}: min={df[col].min()}, max={df[col].max()}, "
                        f"mean={df[col].mean():.2f}"
                    )
            
            # Data (limit per sheet)
            display_df = df.head(150) if len(df) > 150 else df
            lines.append("\n--- Data ---")
            lines.append(display_df.to_string(index=False))
            
            if len(df) > 150:
                lines.append(f"\n... and {len(df) - 150} more rows (truncated)")
            
            all_content.append("\n".join(lines))
        
        content = "\n\n".join(all_content)
        ext = Path(file_path).suffix.lower().lstrip(".")
        
        return ParsedDocument(
            filename=Path(file_path).name,
            file_type=ext,
            content=content,
            metadata={
                "sheet_count": len(xls.sheet_names),
                "sheet_names": xls.sheet_names,
                "total_rows": total_rows,
            },
        )
    except Exception as e:
        raise RuntimeError(f"Failed to parse Excel file: {e}") from e


def parse_json(file_path: str) -> ParsedDocument:
    """Extract text from a JSON file."""
    with open(file_path, "r", encoding="utf-8", errors="replace") as f:
        data = json.load(f)
    
    # Pretty-print JSON for readability
    content = json.dumps(data, indent=2, ensure_ascii=False, default=str)
    
    # Add summary header
    if isinstance(data, list):
        header = f"JSON Array with {len(data)} items.\n\n"
    elif isinstance(data, dict):
        header = f"JSON Object with {len(data)} top-level keys: {', '.join(list(data.keys())[:20])}\n\n"
    else:
        header = ""
    
    return ParsedDocument(
        filename=Path(file_path).name,
        file_type="json",
        content=header + content,
        metadata={"type": type(data).__name__},
    )


def parse_html(file_path: str) -> ParsedDocument:
    """Extract text from an HTML file, stripping tags."""
    with open(file_path, "r", encoding="utf-8", errors="replace") as f:
        raw = f.read()
    
    # Simple tag stripping (no external dependency needed)
    text = re.sub(r'<script[^>]*>.*?</script>', '', raw, flags=re.DOTALL | re.IGNORECASE)
    text = re.sub(r'<style[^>]*>.*?</style>', '', text, flags=re.DOTALL | re.IGNORECASE)
    text = re.sub(r'<[^>]+>', ' ', text)
    text = re.sub(r'\s+', ' ', text).strip()
    
    # Also keep raw for reference
    content = f"--- Extracted Text ---\n{text}\n\n--- Raw HTML (first 5000 chars) ---\n{raw[:5000]}"
    
    return ParsedDocument(
        filename=Path(file_path).name,
        file_type="html",
        content=content,
        metadata={"raw_length": len(raw)},
    )


def parse_pptx(file_path: str) -> ParsedDocument:
    """Extract text from PowerPoint (.pptx) files."""
    try:
        from pptx import Presentation
        prs = Presentation(file_path)
        slides_text = []
        
        for i, slide in enumerate(prs.slides):
            texts = []
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for paragraph in shape.text_frame.paragraphs:
                        text = paragraph.text.strip()
                        if text:
                            texts.append(text)
            if texts:
                slides_text.append(f"[Slide {i + 1}]\n" + "\n".join(texts))
        
        content = "\n\n".join(slides_text)
        return ParsedDocument(
            filename=Path(file_path).name,
            file_type="pptx",
            content=content,
            page_count=len(prs.slides),
            metadata={"slide_count": len(prs.slides)},
        )
    except ImportError:
        # Fall back to text extraction if python-pptx not installed
        return parse_text_fallback(file_path)
    except Exception as e:
        raise RuntimeError(f"Failed to parse PPTX: {e}") from e


def parse_txt(file_path: str) -> ParsedDocument:
    """Extract text from a plain text file."""
    with open(file_path, "r", encoding="utf-8", errors="replace") as f:
        content = f.read()

    return ParsedDocument(
        filename=Path(file_path).name,
        file_type="txt",
        content=content,
        metadata={},
    )


def parse_text_fallback(file_path: str) -> ParsedDocument:
    """
    Universal fallback: try to read any file as text.
    Works for code files, config files, logs, markdown, etc.
    """
    ext = Path(file_path).suffix.lower().lstrip(".")
    try:
        with open(file_path, "r", encoding="utf-8", errors="replace") as f:
            content = f.read()
        
        # Detect if it's actually binary garbage
        # If more than 10% of chars are non-printable, it's likely binary
        non_printable = sum(1 for c in content[:2000] if not c.isprintable() and c not in '\n\r\t')
        if len(content[:2000]) > 0 and (non_printable / min(len(content), 2000)) > 0.10:
            return ParsedDocument(
                filename=Path(file_path).name,
                file_type=ext or "binary",
                content=f"[Binary file detected: {Path(file_path).name}. Unable to extract text content. "
                        f"File size: {os.path.getsize(file_path)} bytes]",
                metadata={"binary": True, "size_bytes": os.path.getsize(file_path)},
            )
        
        return ParsedDocument(
            filename=Path(file_path).name,
            file_type=ext or "text",
            content=content,
            metadata={"parser": "text_fallback"},
        )
    except Exception:
        # Truly unreadable — return a stub
        return ParsedDocument(
            filename=Path(file_path).name,
            file_type=ext or "unknown",
            content=f"[Could not extract text from: {Path(file_path).name}. "
                    f"File size: {os.path.getsize(file_path)} bytes]",
            metadata={"unreadable": True, "size_bytes": os.path.getsize(file_path)},
        )


# ═══════════════════════════════════════════════════════════════════════
#  DISPATCHER
# ═══════════════════════════════════════════════════════════════════════

# Specialized parsers for known formats
PARSERS = {
    ".pdf": parse_pdf,
    ".docx": parse_docx,
    ".csv": parse_csv,
    ".xlsx": parse_excel,
    ".xls": parse_excel,
    ".json": parse_json,
    ".html": parse_html,
    ".htm": parse_html,
    ".pptx": parse_pptx,
    ".txt": parse_txt,
}

# Extensions that are definitely text-based (use text fallback)
TEXT_EXTENSIONS = {
    # Markup & Config
    ".md", ".markdown", ".rst", ".yaml", ".yml", ".toml", ".ini", ".cfg", ".conf",
    ".xml", ".svg", ".env", ".properties",
    # Code
    ".py", ".js", ".ts", ".jsx", ".tsx", ".java", ".c", ".cpp", ".h", ".hpp",
    ".cs", ".go", ".rs", ".rb", ".php", ".swift", ".kt", ".scala", ".r",
    ".sql", ".sh", ".bash", ".zsh", ".ps1", ".bat", ".cmd",
    # Web
    ".css", ".scss", ".sass", ".less",
    # Data
    ".tsv", ".jsonl", ".ndjson",
    # Docs
    ".log", ".tex", ".bib", ".rtf",
}


def parse_file(file_path: str) -> ParsedDocument:
    """
    Parse any file. Uses specialized parsers for known formats,
    falls back to text extraction for everything else.
    Never raises ValueError for unsupported types — always tries its best.
    """
    ext = Path(file_path).suffix.lower()
    
    # 1. Try specialized parser first
    parser = PARSERS.get(ext)
    if parser is not None:
        try:
            return parser(file_path)
        except Exception as e:
            # If specialized parser fails, fall back to text
            try:
                return parse_text_fallback(file_path)
            except Exception:
                return ParsedDocument(
                    filename=Path(file_path).name,
                    file_type=ext.lstrip(".") or "unknown",
                    content=f"[Error parsing {Path(file_path).name}: {str(e)}]",
                    metadata={"error": str(e)},
                )
    
    # 2. Known text extensions → text fallback
    if ext in TEXT_EXTENSIONS:
        return parse_text_fallback(file_path)
    
    # 3. Unknown extension → still try text fallback (gracefully handles binary)
    return parse_text_fallback(file_path)
